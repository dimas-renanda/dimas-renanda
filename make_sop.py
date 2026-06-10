from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette - fashion/modern
COLOR_PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
COLOR_ACCENT    = RGBColor(0xE9, 0x4F, 0x8B)   # pink
COLOR_ACCENT2   = RGBColor(0xFF, 0xC2, 0xD4)   # soft pink
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xF8, 0xF0, 0xF4)   # off-white
COLOR_DARK      = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_MUTED     = RGBColor(0x88, 0x88, 0x99)
COLOR_ECOM      = RGBColor(0xE9, 0x4F, 0x8B)   # pink for ecommerce
COLOR_STOK      = RGBColor(0x6C, 0x5C, 0xE7)   # purple for stok

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p

# ─────────────────────────────────────────────
# SLIDE 1: Cover
# ─────────────────────────────────────────────
slide = blank_slide(prs)

# Full background
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_PRIMARY)

# Decorative accent bar left
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ACCENT)

# Pink decorative blob top-right
bg_rect(slide, Inches(9.5), 0, Inches(3.83), Inches(3), COLOR_ACCENT2)

# White card center
card_x, card_y = Inches(1.2), Inches(1.8)
card_w, card_h = Inches(9), Inches(4.2)
card = bg_rect(slide, card_x, card_y, card_w, card_h, COLOR_WHITE)

# Title
tb = slide.shapes.add_textbox(card_x + Inches(0.4), card_y + Inches(0.35), card_w - Inches(0.8), Inches(1))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "SOP INTERN"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = COLOR_PRIMARY

# Subtitle
tb2 = slide.shapes.add_textbox(card_x + Inches(0.4), card_y + Inches(1.3), card_w - Inches(0.8), Inches(0.6))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "E-Commerce & Stok · Fashion Retail"
run2.font.size = Pt(20)
run2.font.bold = False
run2.font.color.rgb = COLOR_ACCENT

# Divider line
div = slide.shapes.add_shape(1, card_x + Inches(0.4), card_y + Inches(2.05), Inches(6), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = COLOR_ACCENT2
div.line.fill.background()

# Description
tb3 = slide.shapes.add_textbox(card_x + Inches(0.4), card_y + Inches(2.25), card_w - Inches(0.8), Inches(1.5))
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = "Panduan kerja harian intern mencakup operasional e-commerce\ndan pengelolaan stok & packing produk."
run3.font.size = Pt(15)
run3.font.color.rgb = COLOR_DARK

# Bottom label
add_text_box(slide, "Dokumen Internal · Untuk Intern Baru", Inches(1.2), Inches(6.7),
             Inches(9), Inches(0.5), font_size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2: Daftar Isi / Overview
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ACCENT)

# Header bar
bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_PRIMARY)
add_text_box(slide, "RUANG LINGKUP KERJA", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
             font_size=28, bold=True, color=COLOR_WHITE)

# Two columns
col1_x = Inches(0.7)
col2_x = Inches(6.9)
col_y  = Inches(1.8)
col_w  = Inches(5.8)
col_h  = Inches(5)

# Card 1 - Ecommerce
c1 = bg_rect(slide, col1_x, col_y, col_w, col_h, COLOR_WHITE)

tb = slide.shapes.add_textbox(col1_x + Inches(0.25), col_y + Inches(0.2), col_w - Inches(0.5), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run(); run.text = "🛍️  E-COMMERCE"
run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = COLOR_ECOM

items1 = [
    "Flash sale campaign (mingguan)",
    "Chat broadcast (1x seminggu)",
    "Optimasi produk Shopee (tiap 4 jam)",
    "Balas chat termasuk weekend",
    "Atur pengiriman Jumat & Sabtu (hari Minggu)",
]
tb2 = slide.shapes.add_textbox(col1_x + Inches(0.25), col_y + Inches(0.95), col_w - Inches(0.5), col_h - Inches(1.1))
tf2 = tb2.text_frame
tf2.word_wrap = True
for i, item in enumerate(items1):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = f"  ✦  {item}"
    run.font.size = Pt(13.5)
    run.font.color.rgb = COLOR_DARK

# Card 2 - Stok & Packing
c2 = bg_rect(slide, col2_x, col_y, col_w, col_h, COLOR_WHITE)

tb = slide.shapes.add_textbox(col2_x + Inches(0.25), col_y + Inches(0.2), col_w - Inches(0.5), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run(); run.text = "📦  STOK & PACKING"
run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = COLOR_STOK

items2 = [
    "Packing orderan",
    "Stock opname bulanan (hari kerja terakhir)",
    "Cek barang masuk vs surat jalan",
    "QC barang baru (lapor ke produksi)",
    "Rapikan area packing (konten-able)",
]
tb2 = slide.shapes.add_textbox(col2_x + Inches(0.25), col_y + Inches(0.95), col_w - Inches(0.5), col_h - Inches(1.1))
tf2 = tb2.text_frame
tf2.word_wrap = True
for i, item in enumerate(items2):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = f"  ✦  {item}"
    run.font.size = Pt(13.5)
    run.font.color.rgb = COLOR_DARK

# ─────────────────────────────────────────────
# SLIDE 3: SOP E-Commerce (1) - Flash Sale & Broadcast
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ECOM)

# Header
bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_ECOM)
add_text_box(slide, "🛍️  E-COMMERCE  —  Flash Sale & Chat Broadcast", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

def sop_card(slide, x, y, w, h, title, freq, steps, color_accent):
    card = bg_rect(slide, x, y, w, h, COLOR_WHITE)
    # top accent stripe
    stripe = bg_rect(slide, x, y, w, Inches(0.08), color_accent)

    # title
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.18), w - Inches(0.4), Inches(0.55))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = COLOR_PRIMARY

    # freq badge
    tb2 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.72), w - Inches(0.4), Inches(0.35))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run(); run2.text = f"⏱  {freq}"
    run2.font.size = Pt(11.5); run2.font.italic = True; run2.font.color.rgb = color_accent

    # steps
    tb3 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.12), w - Inches(0.4), h - Inches(1.3))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.space_before = Pt(3.5)
        run = p.add_run(); run.text = step
        run.font.size = Pt(12.5); run.font.color.rgb = COLOR_DARK

# Flash Sale card (left, tall)
sop_card(slide,
    x=Inches(0.7), y=Inches(1.65), w=Inches(5.8), h=Inches(5.5),
    title="1. Flash Sale Campaign",
    freq="Mingguan (buat di awal minggu)",
    steps=[
        "1️⃣  Tentukan produk yang akan di-flash sale.",
        "2️⃣  Buat visual/banner promo (koordinasi markom jika perlu).",
        "3️⃣  Input campaign di dashboard Shopee → Flash Sale.",
        "4️⃣  Set harga, stok, dan durasi promo.",
        "5️⃣  Pastikan campaign aktif sebelum promo mulai.",
        "6️⃣  Monitor performa selama campaign berjalan.",
        "7️⃣  Catat hasil (produk terjual, revenue) untuk laporan.",
    ],
    color_accent=COLOR_ECOM
)

# Chat Broadcast card (right)
sop_card(slide,
    x=Inches(6.85), y=Inches(1.65), w=Inches(6.1), h=Inches(5.5),
    title="2. Chat Broadcast",
    freq="1x seminggu (jadwal tetap)",
    steps=[
        "1️⃣  Siapkan draft materi broadcast.",
        "2️⃣  Ajukan ke kakak/tim markom untuk approval.",
        "3️⃣  Setelah disetujui, kirim via fitur broadcast Shopee/WA.",
        "4️⃣  Pantau balasan & follow-up yang masuk.",
        "⚠️  Jangan kirim sebelum materi diapprove.",
    ],
    color_accent=COLOR_ECOM
)

# ─────────────────────────────────────────────
# SLIDE 4: SOP E-Commerce (2) - Optimasi & Balas Chat
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ECOM)

bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_ECOM)
add_text_box(slide, "🛍️  E-COMMERCE  —  Optimasi Produk & Balas Chat", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

sop_card(slide,
    x=Inches(0.7), y=Inches(1.65), w=Inches(5.8), h=Inches(5.5),
    title="3. Optimasi Fitur Shopee",
    freq="Setiap 4 jam sekali (jam kerja)",
    steps=[
        "1️⃣  Buka dashboard Shopee Seller.",
        "2️⃣  Gunakan fitur 'Naikkan Produk' (Boost Product).",
        "3️⃣  Pilih produk prioritas / best seller.",
        "4️⃣  Aktifkan boost — lakukan tiap 4 jam.",
        "5️⃣  Pantau posisi produk di hasil pencarian.",
        "6️⃣  Catat produk mana yang sering di-boost untuk evaluasi.",
    ],
    color_accent=COLOR_ECOM
)

sop_card(slide,
    x=Inches(6.85), y=Inches(1.65), w=Inches(6.1), h=Inches(5.5),
    title="4. Balas Chat Pelanggan",
    freq="Setiap hari termasuk weekend",
    steps=[
        "1️⃣  Cek notifikasi chat Shopee minimal 3x sehari.",
        "2️⃣  Balas semua pertanyaan dengan ramah & informatif.",
        "3️⃣  Jika ada komplain → eskalasi ke kakak (jangan diabaikan).",
        "4️⃣  Gunakan template jawaban untuk FAQ.",
        "5️⃣  Pastikan response rate tetap tinggi (target > 90%).",
        "⚠️  Weekend WAJIB tetap balas — jadwalkan waktu khusus.",
    ],
    color_accent=COLOR_ECOM
)

# ─────────────────────────────────────────────
# SLIDE 5: SOP E-Commerce (3) - Pengiriman Weekend
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ECOM)

bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_ECOM)
add_text_box(slide, "🛍️  E-COMMERCE  —  Pengiriman Weekend", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

# Big centered card
card_x, card_y = Inches(0.7), Inches(1.65)
card_w, card_h = Inches(12.28), Inches(5.5)
bg_rect(slide, card_x, card_y, card_w, card_h, COLOR_WHITE)
bg_rect(slide, card_x, card_y, card_w, Inches(0.08), COLOR_ECOM)

add_text_box(slide, "5. Atur Pengiriman Jumat & Sabtu", card_x + Inches(0.3), card_y + Inches(0.18),
             card_w - Inches(0.6), Inches(0.55), font_size=18, bold=True, color=COLOR_PRIMARY)

tb_freq = slide.shapes.add_textbox(card_x + Inches(0.3), card_y + Inches(0.75), card_w - Inches(0.6), Inches(0.35))
tf_freq = tb_freq.text_frame
p = tf_freq.paragraphs[0]
run = p.add_run(); run.text = "⏱  Setiap hari Minggu — WAJIB dilakukan"
run.font.size = Pt(13); run.font.italic = True; run.font.color.rgb = COLOR_ECOM

# Timeline visual (3 boxes)
boxes = [
    ("JUM'AT", "Order masuk hari Jumat\ndicatat & disiapkan\nuntuk pengiriman"),
    ("SABTU", "Order masuk hari Sabtu\ndicatat & disiapkan\nuntuk pengiriman"),
    ("MINGGU ✅", "WAJIB:\nProses & atur semua\npengiriman Jumat + Sabtu\ndi hari ini"),
]
colors_box = [COLOR_ACCENT2, COLOR_ACCENT2, COLOR_ECOM]
text_colors = [COLOR_DARK, COLOR_DARK, COLOR_WHITE]

for i, (label, desc) in enumerate(boxes):
    bx = card_x + Inches(0.5) + i * Inches(4.0)
    by = card_y + Inches(1.3)
    bw, bh = Inches(3.5), Inches(3.5)
    bg_rect(slide, bx, by, bw, bh, colors_box[i])

    # Label
    tb_l = slide.shapes.add_textbox(bx + Inches(0.15), by + Inches(0.2), bw - Inches(0.3), Inches(0.55))
    tf_l = tb_l.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.alignment = PP_ALIGN.CENTER
    run_l = p_l.add_run(); run_l.text = label
    run_l.font.size = Pt(17); run_l.font.bold = True; run_l.font.color.rgb = text_colors[i]

    # Desc
    tb_d = slide.shapes.add_textbox(bx + Inches(0.15), by + Inches(0.9), bw - Inches(0.3), Inches(2.4))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]
    p_d.alignment = PP_ALIGN.CENTER
    run_d = p_d.add_run(); run_d.text = desc
    run_d.font.size = Pt(13); run_d.font.color.rgb = text_colors[i]

    # Arrow between boxes
    if i < 2:
        arr_x = bx + bw + Inches(0.05)
        arr_y = by + Inches(1.5)
        tb_a = slide.shapes.add_textbox(arr_x, arr_y, Inches(0.4), Inches(0.5))
        tf_a = tb_a.text_frame
        p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        run_a = p_a.add_run(); run_a.text = "→"
        run_a.font.size = Pt(28); run_a.font.color.rgb = COLOR_MUTED

# ─────────────────────────────────────────────
# SLIDE 6: SOP Stok & Packing (1) - Packing & Stock Opname
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_STOK)

bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_STOK)
add_text_box(slide, "📦  STOK & PACKING  —  Packing & Stock Opname", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

sop_card(slide,
    x=Inches(0.7), y=Inches(1.65), w=Inches(5.8), h=Inches(5.5),
    title="1. Packing Orderan",
    freq="Setiap hari (sesuai orderan masuk)",
    steps=[
        "1️⃣  Cek list orderan yang harus di-packing hari ini.",
        "2️⃣  Siapkan produk sesuai orderan (cek SKU & ukuran).",
        "3️⃣  Packing dengan rapi & aman (bubble wrap, dll).",
        "4️⃣  Tempel label pengiriman di luar paket.",
        "5️⃣  Susun paket siap kirim di area yang ditentukan.",
        "6️⃣  Rapikan area packing setelah selesai.",
    ],
    color_accent=COLOR_STOK
)

sop_card(slide,
    x=Inches(6.85), y=Inches(1.65), w=Inches(6.1), h=Inches(5.5),
    title="2. Stock Opname Bulanan",
    freq="1x/bulan — hari kerja terakhir di minggu terakhir",
    steps=[
        "1️⃣  Lakukan bersama kakak (wajib koordinasi jadwal).",
        "2️⃣  Hitung fisik semua stok produk per SKU.",
        "3️⃣  Cocokkan dengan data sistem/catatan stok.",
        "4️⃣  Catat selisih jika ada (lebih/kurang).",
        "5️⃣  Laporkan hasil stock opname ke kakak.",
        "6️⃣  Update data stok setelah opname selesai.",
        "⚠️  Jadwal ditentukan bersama — jangan dilakukan sendiri.",
    ],
    color_accent=COLOR_STOK
)

# ─────────────────────────────────────────────
# SLIDE 7: SOP Stok & Packing (2) - Cek Barang & QC
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_STOK)

bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_STOK)
add_text_box(slide, "📦  STOK & PACKING  —  Cek Barang Masuk & QC", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

sop_card(slide,
    x=Inches(0.7), y=Inches(1.65), w=Inches(5.8), h=Inches(5.5),
    title="3. Cek Barang Masuk vs Surat Jalan",
    freq="Setiap ada pengiriman barang masuk",
    steps=[
        "1️⃣  Terima barang dari kurir/supplier.",
        "2️⃣  Ambil surat jalan / dokumen pengiriman.",
        "3️⃣  Hitung jumlah fisik barang yang datang.",
        "4️⃣  Cocokkan dengan isi surat jalan (qty, item).",
        "5️⃣  Jika sesuai → tandatangani & simpan surat jalan.",
        "6️⃣  Jika ada selisih → JANGAN tanda tangani, lapor ke kakak segera.",
    ],
    color_accent=COLOR_STOK
)

sop_card(slide,
    x=Inches(6.85), y=Inches(1.65), w=Inches(6.1), h=Inches(5.5),
    title="4. QC Barang Baru",
    freq="Langsung setelah barang dicek & diterima",
    steps=[
        "1️⃣  Buka semua paket barang yang baru datang.",
        "2️⃣  Periksa satu per satu: warna, ukuran, kondisi jahitan.",
        "3️⃣  Pisahkan barang OK dan barang bermasalah.",
        "4️⃣  Dokumentasi barang bermasalah (foto).",
        "5️⃣  Laporkan ke kakak → kakak akan follow-up ke tim produksi.",
        "⚠️  QC harus dilakukan segera agar masalah cepat diatasi.",
    ],
    color_accent=COLOR_STOK
)

# ─────────────────────────────────────────────
# SLIDE 8: SOP Stok & Packing (3) - Rapikan Area
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_LIGHT)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_STOK)

bg_rect(slide, Inches(0.35), 0, prs.slide_width - Inches(0.35), Inches(1.4), COLOR_STOK)
add_text_box(slide, "📦  STOK & PACKING  —  Kebersihan & Kerapian Area", Inches(0.7), Inches(0.35),
             Inches(11.5), Inches(0.8), font_size=24, bold=True, color=COLOR_WHITE)

# Full width card
card_x, card_y = Inches(0.7), Inches(1.65)
card_w, card_h = Inches(12.28), Inches(5.5)
bg_rect(slide, card_x, card_y, card_w, card_h, COLOR_WHITE)
bg_rect(slide, card_x, card_y, card_w, Inches(0.08), COLOR_STOK)

add_text_box(slide, "5. Rapikan Area Packing (Tetap Konten-able!)", card_x + Inches(0.3), card_y + Inches(0.18),
             card_w - Inches(0.6), Inches(0.55), font_size=18, bold=True, color=COLOR_PRIMARY)
tb_freq = slide.shapes.add_textbox(card_x + Inches(0.3), card_y + Inches(0.75), card_w - Inches(0.6), Inches(0.35))
tf_freq = tb_freq.text_frame
p = tf_freq.paragraphs[0]
run = p.add_run(); run.text = "⏱  Setelah setiap sesi packing selesai"
run.font.size = Pt(13); run.font.italic = True; run.font.color.rgb = COLOR_STOK

# Two sub-sections
# Left: steps
steps_left = [
    "1️⃣  Bersihkan sisa kardus, plastik, dan material packing.",
    "2️⃣  Kembalikan perlengkapan packing ke tempat semula.",
    "3️⃣  Susun produk yang belum di-packing dengan rapi.",
    "4️⃣  Pastikan label & sticker tersimpan di tempat yang benar.",
    "5️⃣  Cek lantai & meja — bebas dari sampah atau serpihan.",
]
tb_l = slide.shapes.add_textbox(card_x + Inches(0.3), card_y + Inches(1.25), Inches(5.5), Inches(3.8))
tf_l = tb_l.text_frame
tf_l.word_wrap = True
for i, s in enumerate(steps_left):
    p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run(); run.text = s
    run.font.size = Pt(13.5); run.font.color.rgb = COLOR_DARK

# Divider
div = bg_rect(slide, card_x + Inches(6.1), card_y + Inches(1.2), Inches(0.04), Inches(3.8), COLOR_ACCENT2)

# Right: "konten-able" tips
add_text_box(slide, "✨  Standar Konten-able:", card_x + Inches(6.4), card_y + Inches(1.25),
             Inches(5.5), Inches(0.5), font_size=14, bold=True, color=COLOR_STOK)
tips = [
    "📸  Area packing = background konten potensial",
    "🎀  Tatanan produk harus estetik & bersih",
    "🚫  Tidak ada sampah/kardus berantakan di frame",
    "💡  Pencahayaan area harus bersih & tidak berantakan",
    "📦  Produk di-display rapi, tidak menumpuk sembarangan",
]
tb_r = slide.shapes.add_textbox(card_x + Inches(6.4), card_y + Inches(1.85), Inches(5.5), Inches(3.2))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
for i, t in enumerate(tips):
    p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run(); run.text = t
    run.font.size = Pt(13); run.font.color.rgb = COLOR_DARK

# ─────────────────────────────────────────────
# SLIDE 9: Closing / Summary
# ─────────────────────────────────────────────
slide = blank_slide(prs)
bg_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_PRIMARY)
bg_rect(slide, 0, 0, Inches(0.35), prs.slide_height, COLOR_ACCENT)

# Decorative
bg_rect(slide, Inches(9.5), Inches(4.5), Inches(3.83), Inches(3), COLOR_ACCENT2)

# Card
card_x, card_y = Inches(1.2), Inches(1.5)
card_w, card_h = Inches(10.5), Inches(4.8)
bg_rect(slide, card_x, card_y, card_w, card_h, COLOR_WHITE)

add_text_box(slide, "Semangat & Selamat Kerja! 🎉", card_x + Inches(0.4), card_y + Inches(0.35),
             card_w - Inches(0.8), Inches(0.9), font_size=32, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

div = bg_rect(slide, card_x + Inches(2), card_y + Inches(1.35), Inches(6.5), Inches(0.04), COLOR_ACCENT)

reminders = [
    "✅  Selalu tanya kalau bingung — jangan tebak-tebakan",
    "✅  Broadcast & materi promo WAJIB approval dulu",
    "✅  Pengiriman weekend (Minggu) adalah tanggung jawab penuh intern",
    "✅  QC barang segera setelah tiba — jangan ditunda",
    "✅  Jaga kebersihan area packing setiap saat",
]
tb = slide.shapes.add_textbox(card_x + Inches(0.5), card_y + Inches(1.55), card_w - Inches(1), Inches(3))
tf = tb.text_frame
tf.word_wrap = True
for i, r in enumerate(reminders):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run(); run.text = r
    run.font.size = Pt(14); run.font.color.rgb = COLOR_DARK

add_text_box(slide, "Ada pertanyaan? Jangan ragu untuk bertanya! 💬",
             Inches(1.2), Inches(6.55), Inches(10.5), Inches(0.6),
             font_size=13, color=COLOR_MUTED, align=PP_ALIGN.CENTER, italic=True)

# Save
out_path = "/Users/user/.openclaw/workspace/SOP_Intern_Ecommerce_Stok.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
